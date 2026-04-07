"""
Generate PowerPoint Presentation
==================================
Creates a full professional presentation covering:
  - Project overview
  - API research findings (Facebook Ads, Google Ads, YouTube Ads, TikTok)
  - Automation flow design
  - System architecture
  - Implementation roadmap

Usage:
    python scripts/generate_presentation.py

Output:
    Social_Media_Automation_Platform.pptx  (in the project root)
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)   # Microsoft blue
ACCENT_CYAN = RGBColor(0x00, 0xB4, 0xD8)   # Bright cyan
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)
GREEN       = RGBColor(0x00, 0xC8, 0x76)
YELLOW      = RGBColor(0xFF, 0xC1, 0x07)
RED         = RGBColor(0xFF, 0x4D, 0x4D)
CARD_BG     = RGBColor(0x1A, 0x2B, 0x3C)   # Card background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def add_background(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items: list[tuple[str, int]], left, top, width, height,
                   base_font_size=16, color=WHITE, heading_color=ACCENT_CYAN):
    """
    items: list of (text, indent_level) where indent_level 0 = heading, 1 = bullet
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        if level == 0:
            p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(base_font_size - level * 1.5)
        run.font.bold = (level == 0)
        run.font.color.rgb = heading_color if level == 0 else color
        if level >= 1:
            p.space_before = Pt(2)


def add_divider(slide, top, color=ACCENT_BLUE, thickness=Pt(2)):
    line = slide.shapes.add_connector(
        1,  # straight line
        Inches(0.5), top,
        Inches(12.83), top,
    )
    line.line.color.rgb = color
    line.line.width = thickness


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    """Slide 1 — Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, DARK_BG)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_BLUE)

    # Gradient stripe
    add_rect(slide, Inches(0.12), Inches(2.8), SLIDE_W, Inches(0.06), fill_color=ACCENT_CYAN)

    add_text_box(slide, "Social Media Automation Platform",
                 Inches(0.6), Inches(1.0), Inches(12), Inches(1.2),
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text_box(slide, "API Research, Automation Flow & System Architecture",
                 Inches(0.6), Inches(2.2), Inches(12), Inches(0.7),
                 font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)

    add_text_box(slide, "Powered by Claude Code + MCP Tools",
                 Inches(0.6), Inches(3.1), Inches(10), Inches(0.5),
                 font_size=16, color=LIGHT_GREY, align=PP_ALIGN.LEFT, italic=True)

    add_text_box(slide, "March 2026",
                 Inches(0.6), Inches(6.7), Inches(4), Inches(0.4),
                 font_size=13, color=LIGHT_GREY)

    # Platform tags
    platforms = ["Twitter/X", "Instagram", "Facebook", "LinkedIn", "YouTube", "TikTok"]
    for i, p in enumerate(platforms):
        add_rect(slide, Inches(0.6 + i * 2.1), Inches(6.2), Inches(1.9), Inches(0.4),
                 fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(1))
        add_text_box(slide, p, Inches(0.62 + i * 2.1), Inches(6.22), Inches(1.86), Inches(0.36),
                     font_size=11, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    """Slide 2 — Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_BLUE)

    add_text_box(slide, "Agenda", Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1))

    items = [
        ("1   Project Overview & Goals", 0),
        ("2   What We're Building", 0),
        ("3   Facebook Ads API — Research Findings", 0),
        ("4   Google Ads API — Research Findings", 0),
        ("5   YouTube Ads API — Research Findings", 0),
        ("6   TikTok API — Research Findings", 0),
        ("7   API Feasibility Summary", 0),
        ("8   Automation Flow Design (Claude Code + MCP)", 0),
        ("9   System Architecture", 0),
        ("10  Implementation Roadmap", 0),
    ]
    add_bullet_box(slide, items, Inches(0.8), Inches(1.25), Inches(12), Inches(5.8),
                   base_font_size=18, heading_color=LIGHT_GREY)


def slide_overview(prs):
    """Slide 3 — Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_CYAN)

    add_text_box(slide, "Project Overview", Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1), color=ACCENT_CYAN)

    add_text_box(slide,
                 "A unified AI-driven automation platform that manages social media posting, "
                 "analytics collection, and advertising — operated entirely through Claude Code "
                 "chat in VSCode. No web UI. No manual dashboards.",
                 Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
                 font_size=16, color=LIGHT_GREY)

    # Three columns
    cols = [
        ("Social Media\nPosting", "Post to 6 platforms\nautomatically from chat"),
        ("Analytics\nCollection", "Pull metrics from all\nplatforms in one command"),
        ("Ads\nAutomation", "Manage Google Ads &\nFacebook Ads via MCP tools"),
    ]
    col_colors = [ACCENT_BLUE, ACCENT_CYAN, GREEN]
    for i, (title, desc) in enumerate(cols):
        x = Inches(0.6 + i * 4.2)
        add_rect(slide, x, Inches(2.4), Inches(3.9), Inches(3.5),
                 fill_color=CARD_BG, line_color=col_colors[i], line_width=Pt(1.5))
        add_text_box(slide, title, x + Inches(0.15), Inches(2.6), Inches(3.6), Inches(0.9),
                     font_size=20, bold=True, color=col_colors[i], align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x + Inches(0.15), Inches(3.5), Inches(3.6), Inches(1.5),
                     font_size=15, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


def slide_what_we_build(prs):
    """Slide 4 — What We're Building"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=GREEN)

    add_text_box(slide, "What We're Building", Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1), color=GREEN)

    rows = [
        ("Platform", "Organic Posting", "Analytics", "Ads Automation"),
        ("Twitter / X",   "✓", "Limited (Basic tier+)", "—"),
        ("Instagram",     "✓", "✓",                    "Via Meta Marketing API"),
        ("Facebook",      "✓", "✓",                    "Via Meta Marketing API"),
        ("LinkedIn",      "✓", "✓",                    "Via LinkedIn Ads API"),
        ("YouTube",       "✓ (video upload)", "✓",     "Via Google Ads API"),
        ("TikTok",        "✓ (approval req.)", "Limited", "✓ (approval req.)"),
        ("Google Ads",    "—", "—",                    "✓ Full"),
        ("Facebook Ads",  "—", "—",                    "✓ Full"),
    ]

    row_h = Inches(0.5)
    col_widths = [Inches(2.4), Inches(2.8), Inches(3.2), Inches(4.3)]
    col_x = [Inches(0.3), Inches(2.75), Inches(5.6), Inches(8.85)]
    header_colors = [ACCENT_BLUE, ACCENT_BLUE, ACCENT_BLUE, ACCENT_BLUE]

    for r, row in enumerate(rows):
        y = Inches(1.2) + r * row_h
        bg = CARD_BG if r % 2 == 1 else RGBColor(0x14, 0x22, 0x33)
        if r == 0:
            bg = ACCENT_BLUE
        for c, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
            add_rect(slide, x, y, w, row_h, fill_color=bg)
            fc = WHITE if r == 0 else (GREEN if "✓" in cell and "req" not in cell else (YELLOW if "req" in cell or "Limited" in cell else LIGHT_GREY))
            add_text_box(slide, cell, x + Inches(0.05), y + Inches(0.08),
                         w - Inches(0.1), row_h - Inches(0.1),
                         font_size=12, bold=(r == 0), color=fc, align=PP_ALIGN.CENTER)


def slide_facebook_ads(prs):
    """Slide 5 — Facebook Ads API"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=RGBColor(0x18, 0x77, 0xF2))

    add_text_box(slide, "Facebook Ads — Meta Marketing API",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.6), Inches(0.95), Inches(1.4), Inches(0.32),
             fill_color=GREEN, line_color=None)
    add_text_box(slide, "FULLY AUTOMATABLE", Inches(0.62), Inches(0.97), Inches(1.36), Inches(0.28),
                 font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(1.4), color=RGBColor(0x18, 0x77, 0xF2))

    left_items = [
        ("What Can Be Automated", 0),
        ("  ✓  Create campaigns (6 objectives)", 1),
        ("  ✓  Manage ad sets and individual ads", 1),
        ("  ✓  Pause / enable / delete campaigns", 1),
        ("  ✓  Fetch 70+ performance metrics", 1),
        ("  ✓  Upload creative assets", 1),
        ("  ✓  Define targeting and audiences", 1),
        ("  ✓  Adjust budgets post-launch", 1),
        ("Authentication", 0),
        ("  OAuth 2.0 — System User token preferred", 1),
        ("  System User tokens never expire", 1),
        ("  Advanced access requires App Review", 1),
    ]
    add_bullet_box(slide, left_items, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.5),
                   base_font_size=14)

    right_items = [
        ("API Details", 0),
        ("  Version: v25.0 (Jan 2026)", 1),
        ("  Base URL: graph.facebook.com/v25.0", 1),
        ("  Rate limit: 100 mutations/sec", 1),
        ("Key Restrictions", 0),
        ("  All ads go through standard ad review", 1),
        ("  Attribution windows changed Jan 2026", 1),
        ("  Advantage+ structures required (Q1 2026)", 1),
        ("  special_ad_categories required on create", 1),
        ("Key Metrics Available", 0),
        ("  Impressions, reach, clicks, CTR, CPC", 1),
        ("  Spend, conversions, ROAS, frequency", 1),
        ("  Video view-through rates (25/50/100%)", 1),
    ]
    add_bullet_box(slide, right_items, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5),
                   base_font_size=14)


def slide_google_ads(prs):
    """Slide 6 — Google Ads API"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=RGBColor(0x34, 0xA8, 0x53))

    add_text_box(slide, "Google Ads API",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.6), Inches(0.95), Inches(1.4), Inches(0.32),
             fill_color=GREEN, line_color=None)
    add_text_box(slide, "FULLY AUTOMATABLE", Inches(0.62), Inches(0.97), Inches(1.36), Inches(0.28),
                 font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(1.4), color=RGBColor(0x34, 0xA8, 0x53))

    left_items = [
        ("What Can Be Automated", 0),
        ("  ✓  Create/manage Search, Display, Video, PMax", 1),
        ("  ✓  Ad groups, ads, keywords", 1),
        ("  ✓  Budgets and bid strategies", 1),
        ("  ✓  Custom GAQL reports (SQL-like)", 1),
        ("  ✓  Audience and remarketing lists", 1),
        ("  ✓  Conversion tracking setup", 1),
        ("Access Tiers", 0),
        ("  Explorer — 2,880 ops/day (automatic)", 1),
        ("  Basic — 15,000 ops/day (~2 days)", 1),
        ("  Standard — Unlimited (~10 days)", 1),
    ]
    add_bullet_box(slide, left_items, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.5),
                   base_font_size=14)

    right_items = [
        ("API Details", 0),
        ("  Version: v23.1 (Feb 2026)", 1),
        ("  Query Language: GAQL (SQL-like)", 1),
        ("  Required: Developer Token + OAuth 2.0", 1),
        ("  Cost: FREE at all tiers", 1),
        ("Key Restrictions", 0),
        ("  API v19 sunset Feb 11, 2026 — use v23+", 1),
        ("  Explorer excludes account creation & billing", 1),
        ("  Must apply via Google Ads Manager Account", 1),
        ("Key Metrics Available", 0),
        ("  Impressions, clicks, CTR, avg CPC", 1),
        ("  Conversions, cost, ROAS", 1),
        ("  Quality score, search impression share", 1),
    ]
    add_bullet_box(slide, right_items, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5),
                   base_font_size=14)


def slide_youtube_ads(prs):
    """Slide 7 — YouTube Ads"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=RGBColor(0xFF, 0x00, 0x00))

    add_text_box(slide, "YouTube Ads (via Google Ads API)",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.6), Inches(0.95), Inches(1.4), Inches(0.32),
             fill_color=GREEN, line_color=None)
    add_text_box(slide, "FULLY AUTOMATABLE", Inches(0.62), Inches(0.97), Inches(1.36), Inches(0.28),
                 font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(1.4), color=RGBColor(0xFF, 0x00, 0x00))

    add_text_box(slide,
                 "YouTube Ads are NOT a separate API. They are Video campaigns inside "
                 "the Google Ads API. No additional credentials needed.",
                 Inches(0.6), Inches(1.5), Inches(12.4), Inches(0.6),
                 font_size=15, color=YELLOW)

    left_items = [
        ("What Can Be Automated", 0),
        ("  ✓  Create Video & Demand Gen campaigns", 1),
        ("  ✓  YouTube Shorts inventory (Demand Gen)", 1),
        ("  ✓  Performance Max with YouTube placements", 1),
        ("  ✓  Upload video assets (v23.1+ feature)", 1),
        ("  ✓  Targeting, audiences, bidding, budgets", 1),
        ("  ✓  Full analytics via GAQL", 1),
        ("Same credentials as Google Ads", 0),
        ("  Same OAuth 2.0 + Developer Token", 1),
        ("  Same access tiers (Explorer/Basic/Standard)", 1),
    ]
    add_bullet_box(slide, left_items, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8),
                   base_font_size=14)

    right_items = [
        ("Critical 2025–2026 Changes", 0),
        ("  Video Action Campaigns fully deprecated", 1),
        ("  → Replaced by Demand Gen (Q2 2025)", 1),
        ("  Enhanced Conversions now MANDATORY", 1),
        ("  → Required for PMax YouTube (Nov 2025)", 1),
        ("  Demand Gen includes YouTube Shorts", 1),
        ("  → 70B+ daily views inventory", 1),
        ("  tCPM replaced Maximum CPV bidding", 1),
        ("YouTube Shorts Stats (2026)", 0),
        ("  70 billion+ daily views", 1),
        ("  Key placement for Demand Gen campaigns", 1),
    ]
    add_bullet_box(slide, right_items, Inches(6.8), Inches(2.2), Inches(6.0), Inches(4.8),
                   base_font_size=14)


def slide_tiktok(prs):
    """Slide 8 — TikTok API"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=RGBColor(0x01, 0xF0, 0xD0))

    add_text_box(slide, "TikTok Marketing API",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.6), Inches(0.95), Inches(1.8), Inches(0.32),
             fill_color=YELLOW, line_color=None)
    add_text_box(slide, "POSSIBLE — GATED ACCESS", Inches(0.62), Inches(0.97), Inches(1.76), Inches(0.28),
                 font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(1.4), color=RGBColor(0x01, 0xF0, 0xD0))

    add_rect(slide, Inches(0.6), Inches(1.5), Inches(12.4), Inches(0.55),
             fill_color=RGBColor(0x3A, 0x20, 0x00))
    add_text_box(slide,
                 "⚠  TikTok requires MANUAL APPROVAL. Access is not automatic. "
                 "Review takes 2–7 days and can be denied. Design as optional/deferred.",
                 Inches(0.65), Inches(1.52), Inches(12.2), Inches(0.5),
                 font_size=13, color=YELLOW, bold=True)

    left_items = [
        ("What Can Be Automated (after approval)", 0),
        ("  ✓  Post videos (max 20/day hard limit)", 1),
        ("  ✓  Create and manage ad campaigns", 1),
        ("  ✓  Fetch organic video analytics", 1),
        ("  ✓  Fetch ad performance reports", 1),
        ("  ✓  Campaign, ad group, ad management", 1),
        ("Authentication", 0),
        ("  OAuth 2.0 — tokens expire every 24h", 1),
        ("  Must refresh token daily", 1),
        ("  Data refresh required every 15 days", 1),
    ]
    add_bullet_box(slide, left_items, Inches(0.6), Inches(2.15), Inches(5.8), Inches(5.0),
                   base_font_size=14)

    right_items = [
        ("Hard Limitations", 0),
        ("  20 video posts per day (hard limit)", 1),
        ("  Cannot edit published posts via API", 1),
        ("  No custom thumbnails", 1),
        ("  No direct messaging via API", 1),
        ("  No line breaks in captions", 1),
        ("  Not available in all countries", 1),
        ("Approval Process", 0),
        ("  1. Register on TikTok Developer Portal", 1),
        ("  2. Create app with full documentation", 1),
        ("  3. Submit for manual review (2–7 days)", 1),
        ("  4. Can be denied without clear reason", 1),
        ("  5. Tighter requirements added in 2025", 1),
    ]
    add_bullet_box(slide, right_items, Inches(6.8), Inches(2.15), Inches(6.0), Inches(5.0),
                   base_font_size=14)


def slide_feasibility_table(prs):
    """Slide 9 — Feasibility Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_BLUE)

    add_text_box(slide, "API Feasibility Summary",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1))

    rows = [
        ("Platform",      "Automation", "Create Campaigns", "Analytics", "Auth",              "Friction"),
        ("Meta Ads",      "✓ Full",     "✓ Yes",            "70+ metrics","OAuth 2.0\nSys User","App Review\nrequired"),
        ("Google Ads",    "✓ Full",     "✓ Yes",            "GAQL unlimited","OAuth 2.0\nDev Token","Explorer auto\n~2-10 days"),
        ("YouTube Ads",   "✓ Full",     "✓ Via Google Ads", "GAQL unlimited","Same as\nGoogle Ads","Same as\nGoogle Ads"),
        ("TikTok",        "⚠ Gated",    "✓ After approval","Limited",    "OAuth 2.0\n24h tokens","Manual review\nCan be denied"),
    ]

    col_w = [Inches(2.0), Inches(1.8), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]
    col_x = [Inches(0.3), Inches(2.35), Inches(4.2), Inches(6.45), Inches(8.7), Inches(10.95)]
    row_h = Inches(0.9)

    status_colors = {
        "✓ Full": GREEN, "✓ Yes": GREEN, "✓ Via Google Ads": GREEN,
        "✓ After approval": GREEN, "✓ GAQL unlimited": GREEN,
        "GAQL unlimited": GREEN, "70+ metrics": GREEN, "Limited": YELLOW,
        "⚠ Gated": YELLOW, "App Review\nrequired": YELLOW,
        "Explorer auto\n~2-10 days": YELLOW, "Same as\nGoogle Ads": GREEN,
        "Manual review\nCan be denied": RED,
    }

    for r, row in enumerate(rows):
        y = Inches(1.25) + r * row_h
        for c, (cell, w, x) in enumerate(zip(row, col_w, col_x)):
            bg = ACCENT_BLUE if r == 0 else (CARD_BG if r % 2 == 1 else RGBColor(0x14, 0x22, 0x33))
            add_rect(slide, x, y, w, row_h, fill_color=bg)
            fc = status_colors.get(cell, WHITE) if r > 0 else WHITE
            add_text_box(slide, cell, x + Inches(0.04), y + Inches(0.05),
                         w - Inches(0.08), row_h - Inches(0.1),
                         font_size=12, bold=(r == 0), color=fc, align=PP_ALIGN.CENTER)

    # Legend
    for label, color in [("✓ Supported", GREEN), ("⚠ Restricted", YELLOW), ("✗ Not available", RED)]:
        pass  # Could add legend boxes here; keeping clean


def slide_automation_flow(prs):
    """Slide 10 — Automation Flow Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_CYAN)

    add_text_box(slide, "Automation Flow — Claude Code + MCP",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1), color=ACCENT_CYAN)

    # Flow boxes
    steps = [
        ("1", "User Request", "User types a natural\nlanguage request\nin Claude Code chat"),
        ("2", "Claude Code\nUnderstands", "AI interprets intent,\ngenerates content,\ndecides which tools to call"),
        ("3", "MCP Tool\nCalled", "Specific MCP tool\nis triggered with\nthe right parameters"),
        ("4", "API Executed", "Platform API is called\n(Twitter, Google Ads,\netc.)"),
        ("5", "Result\nReported", "Claude Code reports\nresult back in\nthe chat"),
    ]

    step_colors = [ACCENT_BLUE, ACCENT_CYAN, GREEN, YELLOW, ACCENT_BLUE]
    box_w = Inches(2.2)
    arrow_w = Inches(0.4)
    total = len(steps) * box_w + (len(steps) - 1) * arrow_w
    start_x = (SLIDE_W - total) / 2

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (box_w + arrow_w)
        color = step_colors[i]
        add_rect(slide, x, Inches(1.5), box_w, Inches(3.8),
                 fill_color=CARD_BG, line_color=color, line_width=Pt(2))
        # Step number circle area
        add_rect(slide, x, Inches(1.5), box_w, Inches(0.6), fill_color=color)
        add_text_box(slide, num, x, Inches(1.52), box_w, Inches(0.56),
                     font_size=22, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x + Inches(0.05), Inches(2.2), box_w - Inches(0.1), Inches(0.8),
                     font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x + Inches(0.1), Inches(3.1), box_w - Inches(0.2), Inches(2.0),
                     font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

        # Arrow
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.05)
            add_text_box(slide, "→", ax, Inches(2.8), arrow_w - Inches(0.1), Inches(0.5),
                         font_size=22, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "No UI required — Claude Code chat IS the interface. "
                 "All actions triggered through natural language.",
                 Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
                 font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


def slide_workflow_examples(prs):
    """Slide 11 — Workflow Examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=GREEN)

    add_text_box(slide, "Real-World Workflow Examples",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1), color=GREEN)

    examples = [
        ('"Post our product launch announcement to Twitter and LinkedIn"',
         'Claude generates platform-tailored versions → calls post_to_platform x2 → reports post IDs'),
        ('"Schedule an Instagram post for Friday 9am with this image"',
         'Claude formats content → calls schedule_post with scheduled_time → confirms job ID'),
        ('"How did all our platforms perform last week?"',
         'Calls get_all_analytics → aggregates data from 6 platforms → returns unified report in chat'),
        ('"Show me Google Ads performance this month, pause any campaign over $5 CPC"',
         'Calls list_google_campaigns → get_google_campaign_metrics → pause_google_campaign for matches'),
        ('"Create a Facebook Awareness campaign with $50/day budget"',
         'Calls create_facebook_campaign(objective=OUTCOME_AWARENESS, daily_budget_cents=5000)'),
    ]

    for i, (query, action) in enumerate(examples):
        y = Inches(1.2) + i * Inches(1.12)
        add_rect(slide, Inches(0.3), y, Inches(12.7), Inches(1.0),
                 fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(1))
        add_text_box(slide, f"👤 {query}", Inches(0.45), y + Inches(0.04),
                     Inches(12.3), Inches(0.42), font_size=12, bold=True, color=ACCENT_CYAN)
        add_text_box(slide, f"🤖 {action}", Inches(0.45), y + Inches(0.5),
                     Inches(12.3), Inches(0.42), font_size=11, color=LIGHT_GREY)


def slide_architecture(prs):
    """Slide 12 — System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_BLUE)

    add_text_box(slide, "System Architecture",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1))

    layers = [
        ("Claude Code (VSCode)", "User chat interface · AI brain · Content generation",
         ACCENT_BLUE, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.75)),
        ("MCP Server  (mcp-server/server.py)", "social_tools · ads_tools · 13 tools exposed via stdio",
         ACCENT_CYAN, Inches(0.3), Inches(2.15), Inches(12.7), Inches(0.75)),
    ]

    for title, sub, color, x, y, w, h in layers:
        add_rect(slide, x, y, w, h, fill_color=CARD_BG, line_color=color, line_width=Pt(2))
        add_text_box(slide, title, x + Inches(0.1), y + Inches(0.04),
                     w - Inches(0.2), Inches(0.38), font_size=14, bold=True, color=color)
        add_text_box(slide, sub, x + Inches(0.1), y + Inches(0.42),
                     w - Inches(0.2), Inches(0.3), font_size=11, color=LIGHT_GREY)

    # Arrow
    add_text_box(slide, "↕  MCP Protocol (stdio)", Inches(5.5), Inches(2.0),
                 Inches(2.8), Inches(0.3), font_size=10, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # Integration boxes row
    integrations = [
        ("Twitter/X", ACCENT_BLUE),
        ("Instagram", RGBColor(0xE1, 0x30, 0x6C)),
        ("Facebook", RGBColor(0x18, 0x77, 0xF2)),
        ("LinkedIn", RGBColor(0x00, 0x77, 0xB5)),
        ("YouTube", RGBColor(0xFF, 0x00, 0x00)),
        ("TikTok", RGBColor(0x01, 0xF0, 0xD0)),
        ("Google Ads", RGBColor(0x34, 0xA8, 0x53)),
        ("Facebook Ads", RGBColor(0x18, 0x77, 0xF2)),
    ]
    box_w = Inches(1.5)
    start_x = Inches(0.3)
    for i, (name, color) in enumerate(integrations):
        x = start_x + i * (box_w + Inches(0.08))
        add_rect(slide, x, Inches(3.1), box_w, Inches(0.65),
                 fill_color=CARD_BG, line_color=color, line_width=Pt(1.5))
        add_text_box(slide, name, x + Inches(0.04), Inches(3.15),
                     box_w - Inches(0.08), Inches(0.55),
                     font_size=11, bold=True, color=color, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Platform Integrations (integrations/)", Inches(0.3), Inches(2.98),
                 Inches(4), Inches(0.3), font_size=11, color=LIGHT_GREY)
    add_divider(slide, Inches(3.85), color=DARK_GREY, thickness=Pt(1))

    # Automation layer
    automation = [
        ("Scheduler", "Persistent post queue\nPolls every 60s"),
        ("Analytics\nCollector", "All platforms\nEvery 24 hours"),
        ("Campaign\nMonitor", "Spend/CTR/CPC\nalerts every 60min"),
    ]
    add_text_box(slide, "Automation Layer (automation/)", Inches(0.3), Inches(3.9),
                 Inches(4), Inches(0.3), font_size=11, color=LIGHT_GREY)
    for i, (name, desc) in enumerate(automation):
        x = Inches(0.3 + i * 4.35)
        add_rect(slide, x, Inches(4.25), Inches(4.1), Inches(1.4),
                 fill_color=CARD_BG, line_color=GREEN, line_width=Pt(1))
        add_text_box(slide, name, x + Inches(0.1), Inches(4.3),
                     Inches(3.9), Inches(0.5), font_size=13, bold=True, color=GREEN)
        add_text_box(slide, desc, x + Inches(0.1), Inches(4.8),
                     Inches(3.9), Inches(0.7), font_size=11, color=LIGHT_GREY)

    add_text_box(slide, "Data: data/scheduled_posts.json  ·  data/analytics/  ·  data/alerts/",
                 Inches(0.3), Inches(5.8), Inches(12.7), Inches(0.35),
                 font_size=11, color=DARK_GREY, align=PP_ALIGN.CENTER, italic=True)


def slide_mcp_tools(prs):
    """Slide 13 — MCP Tools Reference"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_CYAN)

    add_text_box(slide, "MCP Tools — Full Reference (13 Tools)",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1), color=ACCENT_CYAN)

    # Social tools column
    add_text_box(slide, "Social Media Tools", Inches(0.3), Inches(1.2), Inches(6.0), Inches(0.4),
                 font_size=15, bold=True, color=ACCENT_CYAN)
    social = [
        ("post_to_platform",     "Post to a single platform immediately"),
        ("post_to_all_platforms","Post to multiple platforms concurrently"),
        ("schedule_post",        "Schedule a post for a future time"),
        ("get_post_analytics",   "Get metrics for a specific post"),
        ("get_account_analytics","Account-level metrics for one platform"),
        ("get_all_analytics",    "Aggregate analytics from all platforms"),
    ]
    for i, (name, desc) in enumerate(social):
        y = Inches(1.65) + i * Inches(0.75)
        add_rect(slide, Inches(0.3), y, Inches(6.0), Inches(0.68),
                 fill_color=CARD_BG, line_color=ACCENT_CYAN, line_width=Pt(1))
        add_text_box(slide, name, Inches(0.4), y + Inches(0.04),
                     Inches(5.8), Inches(0.3), font_size=12, bold=True, color=ACCENT_CYAN)
        add_text_box(slide, desc, Inches(0.4), y + Inches(0.34),
                     Inches(5.8), Inches(0.28), font_size=11, color=LIGHT_GREY)

    # Ads tools column
    add_text_box(slide, "Ads Tools", Inches(6.7), Inches(1.2), Inches(6.0), Inches(0.4),
                 font_size=15, bold=True, color=GREEN)
    ads = [
        ("list_google_campaigns",          "List all Google / YouTube Ads campaigns"),
        ("get_google_campaign_metrics",    "CTR, CPC, spend, conversions (Google)"),
        ("create_google_campaign",         "Create Search / Display / Video campaign"),
        ("pause_google_campaign",          "Pause a Google Ads campaign"),
        ("list_facebook_campaigns",        "List all Facebook Ads campaigns"),
        ("get_facebook_campaign_metrics",  "CTR, CPC, spend, ROAS (Facebook)"),
        ("create_facebook_campaign",       "Create a Facebook Ads campaign"),
        ("pause_facebook_campaign",        "Pause a Facebook Ads campaign"),
        ("get_facebook_ads_account_metrics","Total spend summary across all FB campaigns"),
    ]
    for i, (name, desc) in enumerate(ads):
        y = Inches(1.65) + i * Inches(0.6)
        add_rect(slide, Inches(6.7), y, Inches(6.3), Inches(0.53),
                 fill_color=CARD_BG, line_color=GREEN, line_width=Pt(1))
        add_text_box(slide, name, Inches(6.8), y + Inches(0.02),
                     Inches(6.1), Inches(0.25), font_size=11, bold=True, color=GREEN)
        add_text_box(slide, desc, Inches(6.8), y + Inches(0.27),
                     Inches(6.1), Inches(0.22), font_size=10, color=LIGHT_GREY)


def slide_roadmap(prs):
    """Slide 14 — Implementation Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=YELLOW)

    add_text_box(slide, "Implementation Roadmap",
                 Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.1), color=YELLOW)

    phases = [
        ("Phase 1", "Core Social Posting", GREEN,
         ["MCP server running (done)", "Twitter, Instagram, Facebook, LinkedIn, YouTube (done)",
          "post_to_platform & post_to_all_platforms", "Multi-platform concurrent posting"]),
        ("Phase 2", "Analytics", ACCENT_CYAN,
         ["get_all_analytics — all platforms in one call (done)", "Per-platform account analytics",
          "Per-post analytics", "Analytics snapshots saved to disk"]),
        ("Phase 3", "Facebook Ads", RGBColor(0x18, 0x77, 0xF2),
         ["Meta Marketing API integration (done)", "list/create/pause/metrics tools (done)",
          "System User token setup", "App Review for Advanced access"]),
        ("Phase 4", "Google + YouTube Ads", RGBColor(0x34, 0xA8, 0x53),
         ["Google Ads API integration (done)", "YouTube Ads via Video campaign type",
          "Apply for Basic/Standard tier access", "Enhanced Conversions for YouTube PMax"]),
        ("Phase 5", "TikTok (Optional)", YELLOW,
         ["Apply for TikTok API access", "Wait for manual approval (2–7 days)",
          "Organic posting + Marketing API", "Integrate if approval granted"]),
    ]

    ph_w = Inches(2.4)
    for i, (phase, title, color, items) in enumerate(phases):
        x = Inches(0.3) + i * (ph_w + Inches(0.12))
        add_rect(slide, x, Inches(1.25), ph_w, Inches(5.8),
                 fill_color=CARD_BG, line_color=color, line_width=Pt(2))
        add_rect(slide, x, Inches(1.25), ph_w, Inches(0.55), fill_color=color)
        add_text_box(slide, phase, x, Inches(1.27), ph_w, Inches(0.51),
                     font_size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x + Inches(0.08), Inches(1.88),
                     ph_w - Inches(0.16), Inches(0.5),
                     font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, f"• {item}", x + Inches(0.1), Inches(2.45) + j * Inches(0.75),
                         ph_w - Inches(0.2), Inches(0.65),
                         font_size=10, color=LIGHT_GREY)


def slide_closing(prs):
    """Slide 15 — Closing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT_BLUE)
    add_rect(slide, Inches(0.12), Inches(3.6), SLIDE_W, Inches(0.06), fill_color=ACCENT_CYAN)

    add_text_box(slide, "Key Takeaways",
                 Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.15))

    takeaways = [
        ("✓", "Facebook Ads, Google Ads, and YouTube Ads are all fully automatable via official APIs."),
        ("✓", "TikTok API access is possible but requires manual approval — treat as optional."),
        ("✓", "No UI needed — Claude Code chat IS the interface for all operations."),
        ("✓", "MCP server exposes 13 tools covering all social posting and ad automation."),
        ("✓", "Background automation handles scheduling, analytics collection, and campaign monitoring."),
        ("→", "Next: Configure API keys → Test connections → Register MCP server → Start automating."),
    ]

    for i, (icon, text) in enumerate(takeaways):
        y = Inches(1.3) + i * Inches(0.82)
        color = GREEN if icon == "✓" else YELLOW
        add_text_box(slide, icon, Inches(0.6), y, Inches(0.4), Inches(0.5),
                     font_size=20, bold=True, color=color)
        add_text_box(slide, text, Inches(1.1), y + Inches(0.04), Inches(11.5), Inches(0.55),
                     font_size=15, color=LIGHT_GREY)

    add_text_box(slide, "Powered by Claude Code + MCP  ·  March 2026",
                 Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
                 font_size=13, color=DARK_GREY, align=PP_ALIGN.CENTER, italic=True)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_presentation(output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs);               print("  [OK] Slide 1  - Title")
    slide_agenda(prs);              print("  [OK] Slide 2  - Agenda")
    slide_overview(prs);            print("  [OK] Slide 3  - Project Overview")
    slide_what_we_build(prs);       print("  [OK] Slide 4  - What We're Building")
    slide_facebook_ads(prs);        print("  [OK] Slide 5  - Facebook Ads API")
    slide_google_ads(prs);          print("  [OK] Slide 6  - Google Ads API")
    slide_youtube_ads(prs);         print("  [OK] Slide 7  - YouTube Ads API")
    slide_tiktok(prs);              print("  [OK] Slide 8  - TikTok API")
    slide_feasibility_table(prs);   print("  [OK] Slide 9  - Feasibility Summary")
    slide_automation_flow(prs);     print("  [OK] Slide 10 - Automation Flow")
    slide_workflow_examples(prs);   print("  [OK] Slide 11 - Workflow Examples")
    slide_architecture(prs);        print("  [OK] Slide 12 - System Architecture")
    slide_mcp_tools(prs);           print("  [OK] Slide 13 - MCP Tools Reference")
    slide_roadmap(prs);             print("  [OK] Slide 14 - Implementation Roadmap")
    slide_closing(prs);             print("  [OK] Slide 15 - Key Takeaways")

    prs.save(output_path)
    print(f"\nDone! Presentation saved: {output_path}")


if __name__ == "__main__":
    root = Path(__file__).parent.parent
    out = str(root / "Social_Media_Automation_Platform.pptx")
    build_presentation(out)
